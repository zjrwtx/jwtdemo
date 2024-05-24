from fastapi import FastAPI, File, UploadFile, Form, HTTPException, Depends
from typing import List, Optional
from PyPDF2 import PdfReader
from pptx import Presentation
from io import BytesIO
from fastapi.responses import JSONResponse
from fastapi import HTTPException, status
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy import create_engine, Column, String, Boolean
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from pydantic import BaseModel, EmailStr
from passlib.context import CryptContext
from jose import JWTError, jwt
from datetime import datetime, timedelta, timezone
from dotenv import load_dotenv
import requests
import os
class User(BaseModel):
    username: str
    email: Optional[EmailStr] = None
    disabled: Optional[bool] = None

class UserInDB(User):
    hashed_password: str

class UserCreate(BaseModel):
    username: str
    email: EmailStr
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str

class TokenData(BaseModel):
    username: Optional[str] = None

class UsernameCheck(BaseModel):
    username: str

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Import User class if it's defined in another module
# from models import User

load_dotenv()

DATABASE_URL = os.getenv("DATABASE_URL")
SECRET_KEY = os.getenv("SECRET_KEY")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30
API_KEY = os.getenv('API_KEY')
API_url = os.getenv('API_url')

engine = create_engine(DATABASE_URL)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

class UserModel(Base):
    __tablename__ = "users"
    username = Column(String(50), primary_key=True, index=True)
    email = Column(String(100), unique=True, index=True)
    hashed_password = Column(String(100))
    disabled = Column(Boolean, default=False)

Base.metadata.create_all(bind=engine)

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

token_blacklist = []

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

def verify_password(plain_password, hashed_password):
    return pwd_context.verify(plain_password, hashed_password)

def get_password_hash(password):
    return pwd_context.hash(password)

def get_user(db: Session, username: str):
    return db.query(UserModel).filter(UserModel.username == username).first()

def authenticate_user(db: Session, username: str, password: str):
    user = get_user(db, username)
    if not user or not verify_password(password, user.hashed_password):
        return False
    return user

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + expires_delta if expires_delta else datetime.now(timezone.utc) + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

async def get_current_user(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    if token in token_blacklist:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Token has been revoked", headers={"WWW-Authenticate": "Bearer"})
    
    credentials_exception = HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Could not validate credentials", headers={"WWW-Authenticate": "Bearer"})
    
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise credentials_exception
        token_data = TokenData(username=username)
    except JWTError:
        raise credentials_exception
    
    user = get_user(db, username=token_data.username)
    if user is None:
        raise credentials_exception
    return user

async def get_current_active_user(current_user: User = Depends(get_current_user)):
    if current_user.disabled:
        raise HTTPException(status_code=400, detail="Inactive user")
    return current_user

def maketopics(text, typeoftopic, numoftopic):
    headers = {'Content-Type': 'application/json', 'Authorization': f"Bearer {API_KEY}"}
    json_data = {
        'model': 'yi-large',
        'messages': [
            {'role': 'system', 'content': '你是一个专业的出题者，能对用户输入的内容生成详细的高质量的题目，且带上答案与解释'},
            {'role': 'user', 'content': f"请你把我输入的内容生成{numoftopic}道详细的高质量的{typeoftopic}类型题目，且带上答案与解释: {text}"},
        ],
        'temperature': 0.3,
        "max_tokens": 10000,
        "top_p": 1.0,
    }
    response = requests.post(API_url, headers=headers, json=json_data)
    try:
        return response.json()['choices'][0]['message']['content']
    except (KeyError, IndexError, AttributeError):
        return "出错啦！ 请联系和push一下开发者解决问题"

def mindmap(text):
    headers = {'Content-Type': 'application/json', 'Authorization': f"Bearer {API_KEY}"}
    json_data = {
        'model': 'yi-large',
        'messages': [
            {'role': 'system', 'content': '你是一个专业的思维导图专家，能对用户输入的内容生成详细的高质量的markdown格式的思维导图'},
            {'role': 'user', 'content': f"请你把我输入的内容生成详细的高质量的markdown格式的思维导图{text}"},
        ],
        'temperature': 0.3,
        "max_tokens": 10000,
        "top_p": 1.0,
    }
    response = requests.post(API_url, headers=headers, json=json_data)
    try:
        return response.json()['choices'][0]['message']['content']
    except (KeyError, IndexError, AttributeError):
        return "出错啦！ 请联系和push一下开发者解决问题"

def extract_text_from_pdf(content):
    pdf_file = BytesIO(content)
    pdf_reader = PdfReader(pdf_file)
    text = ''
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text += page.extract_text()
    return text

def extract_text_from_pptx(content):
    pptx_file = BytesIO(content)
    prs = Presentation(pptx_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)


@app.post("/token", response_model=Token)
async def login_for_access_token(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = authenticate_user(db, form_data.username, form_data.password)
    if not user:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Incorrect username or password", headers={"WWW-Authenticate": "Bearer"})
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(data={"sub": user.username}, expires_delta=access_token_expires)
    return {"access_token": access_token, "token_type": "bearer"}

@app.get("/users/me/", response_model=User)
async def read_users_me(current_user: User = Depends(get_current_active_user)):
    return current_user

@app.post("/register", response_model=User)
async def register(user: UserCreate, db: Session = Depends(get_db)):
    db_user = get_user(db, user.username)
    if db_user:
        raise HTTPException(status_code=400, detail="Username already registered")
    hashed_password = get_password_hash(user.password)
    db_user = UserModel(username=user.username, email=user.email, hashed_password=hashed_password)
    db.add(db_user)
    db.commit()
    db.refresh(db_user)
    return db_user

@app.post("/check_username")
async def check_username(user: UsernameCheck, db: Session = Depends(get_db)):
    db_user = get_user(db, user.username)
    if db_user:
        raise HTTPException(status_code=400, detail="Username already taken")
    return {"message": "Username is available"}

@app.post("/topics")
async def upload_files(files: List[UploadFile] = File(...), typeoftopic: str = Form(...), numoftopic: int = Form(...)):
    responses = []
    for file in files:
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()
        if ext not in [".pdf", ".pptx"]:
            raise HTTPException(status_code=400, detail="Unsupported file type. Only PDF and PPTX files are allowed.")
        content = await file.read()
        extracted_text = extract_text_from_pdf(content) if ext == ".pdf" else extract_text_from_pptx(content)
        api_response = maketopics(extracted_text, typeoftopic, numoftopic)
        responses.append({"filename": filename, "api_response": api_response})
    return {"responses": responses}

@app.post("/mindmap")
async def upload_files(files: List[UploadFile] = File(...)):
    responses = []
    for file in files:
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()
        if ext not in [".pdf", ".pptx"]:
            raise HTTPException(status_code=400, detail="Unsupported file type. Only PDF and PPTX files are allowed.")
        content = await file.read()
        extracted_text = extract_text_from_pdf(content) if ext == ".pdf" else extract_text_from_pptx(content)
        api_response = mindmap(extracted_text)
        responses.append({"filename": filename, "api_response": api_response})
    return {"responses": responses}

@app.post("/check_user")
async def check_user(user: UsernameCheck, db: Session = Depends(get_db)):
    db_user = get_user(db, user.username)
    return {"exists": db_user is not None}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
